"""Tests for VLM-gated image extraction in the document loaders.

Covers:
- Embedded image extraction with inline placeholders for DOCX/PPTX/PDF
- Sequential integrity (placeholder sits between before/after text)
- Manifest generation
- The ``extract_images=False`` no-op path (existing behaviour)
- Small/decorative image filtering
- ``build_doc_research`` factory inlining image refs
- ``view_image`` question parameter + generic-description caching
"""

from __future__ import annotations

import io
import os
import re
from pathlib import Path

import pytest
from PIL import Image

from tools.research._loaders import _doc_key, _keep_image, load_source
from tools.research.doc_researcher import build_doc_research


# ---------------------------------------------------------------------------
# Fixtures / helpers
# ---------------------------------------------------------------------------

def _noisy_image(w: int = 160, h: int = 160) -> Image.Image:
    """Random-noise RGB image — compresses poorly so it clears the size floor."""
    return Image.frombytes("RGB", (w, h), os.urandom(w * h * 3))


def _png_bytes(w: int = 160, h: int = 160) -> bytes:
    buf = io.BytesIO()
    _noisy_image(w, h).save(buf, "PNG")
    return buf.getvalue()


def _write_png(path: Path, w: int = 160, h: int = 160) -> Path:
    path.write_bytes(_png_bytes(w, h))
    return path


def _make_docx(path: Path, image_path: Path) -> Path:
    from docx import Document as DocxDocument

    doc = DocxDocument()
    doc.add_paragraph("Intro paragraph before the figure.")
    doc.add_picture(str(image_path))
    doc.add_paragraph("Conclusion paragraph after the figure.")
    doc.save(str(path))
    return path


def _make_pptx(path: Path, image_path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Quarterly results slide"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), Inches(3), Inches(3))
    prs.save(str(path))
    return path


def _make_pdf(path: Path) -> Path:
    # Pillow can author a single-image PDF, which pypdf reads back as an
    # embedded image XObject.
    _noisy_image(240, 240).save(str(path), "PDF", resolution=100.0)
    return path


def _placeholder_ref(content: str) -> str:
    m = re.search(r"\[IMAGE: (\S+) ", content)
    assert m, f"no image placeholder found in:\n{content}"
    return m.group(1)


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

async def test_docx_extracts_image_in_document_order(tmp_path: Path):
    img = _write_png(tmp_path / "fig.png")
    docx = _make_docx(tmp_path / "report.docx", img)

    docs = await load_source(
        "report.docx", files_dir=tmp_path, extract_images=True,
    )
    content = docs[0].page_content

    # Sequential integrity: placeholder sits between before/after text.
    i_before = content.index("Intro paragraph before")
    i_img = content.index("[IMAGE:")
    i_after = content.index("Conclusion paragraph after")
    assert i_before < i_img < i_after

    # Referential integrity: the placeholder path points at a real file.
    ref = _placeholder_ref(content)
    assert ref.startswith("/doc_images/")
    extracted = tmp_path / ref.lstrip("/")
    assert extracted.is_file()

    # Manifest written alongside the extracted image.
    manifest = extracted.parent / "manifest.json"
    assert manifest.is_file()


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

async def test_pptx_extracts_image_within_slide(tmp_path: Path):
    img = _write_png(tmp_path / "chart.png")
    pptx = _make_pptx(tmp_path / "deck.pptx", img)

    docs = await load_source(
        "deck.pptx", files_dir=tmp_path, extract_images=True,
    )
    content = docs[0].page_content

    assert "Quarterly results slide" in content
    ref = _placeholder_ref(content)
    assert "slide 1" in content
    assert (tmp_path / ref.lstrip("/")).is_file()


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

async def test_pdf_extracts_image_with_page_tag(tmp_path: Path):
    pdf = _make_pdf(tmp_path / "scan.pdf")

    docs = await load_source(
        "scan.pdf", files_dir=tmp_path, extract_images=True,
    )
    content = docs[0].page_content

    assert "[IMAGE:" in content
    assert "page 1" in content
    ref = _placeholder_ref(content)
    assert (tmp_path / ref.lstrip("/")).is_file()


# ---------------------------------------------------------------------------
# No-op path
# ---------------------------------------------------------------------------

async def test_extract_images_false_is_noop(tmp_path: Path):
    img = _write_png(tmp_path / "fig.png")
    _make_docx(tmp_path / "report.docx", img)

    docs = await load_source(
        "report.docx", files_dir=tmp_path, extract_images=False,
    )
    content = docs[0].page_content

    assert "[IMAGE:" not in content
    assert "Intro paragraph before" in content
    assert not (tmp_path / "doc_images").exists()


# ---------------------------------------------------------------------------
# Filtering + doc_key
# ---------------------------------------------------------------------------

def test_keep_image_filters_small_and_tiny():
    # Below the byte floor.
    assert _keep_image(b"x" * 100) is False
    # Real image but below the dimension floor.
    buf = io.BytesIO()
    _noisy_image(16, 16).save(buf, "PNG")
    assert _keep_image(buf.getvalue()) is False
    # Real, large enough image is kept.
    assert _keep_image(_png_bytes(160, 160)) is True


def test_doc_key_changes_when_file_edited(tmp_path: Path):
    p = tmp_path / "f.txt"
    p.write_text("one")
    k1 = _doc_key(p)
    # Same content/mtime → stable key (reuses extraction).
    assert _doc_key(p) == k1
    # Editing bumps mtime → new key (fresh extraction).
    os.utime(p, (1, 1))
    assert _doc_key(p) != k1


# ---------------------------------------------------------------------------
# build_doc_research factory
# ---------------------------------------------------------------------------

async def test_build_doc_research_inlines_image_refs(tmp_path: Path):
    img = _write_png(tmp_path / "fig.png")
    _make_docx(tmp_path / "report.docx", img)

    tool = build_doc_research(files_dir=tmp_path, extract_images=True)
    result = await tool.ainvoke({"source": "report.docx", "query": "figure"})

    assert "[IMAGE: /doc_images/" in result


async def test_build_doc_research_default_no_extraction(tmp_path: Path):
    img = _write_png(tmp_path / "fig.png")
    _make_docx(tmp_path / "report.docx", img)

    tool = build_doc_research(files_dir=tmp_path, extract_images=False)
    result = await tool.ainvoke({"source": "report.docx", "query": "figure"})

    assert "[IMAGE:" not in result


# ---------------------------------------------------------------------------
# DocReader preserves exact image paths despite summarisation
# ---------------------------------------------------------------------------

async def test_doc_reader_preserves_image_paths_through_summarisation(tmp_path: Path):
    """The summarising LLM paraphrases content; the exact /doc_images/ paths
    must still reach the agent via the appended reference list."""
    from langchain_core.language_models.fake_chat_models import FakeListChatModel

    from tools.research.doc_reader import DocReader

    img = _write_png(tmp_path / "fig.png")
    _make_docx(tmp_path / "report.docx", img)

    # A paraphrasing summariser that never echoes the literal /doc_images path.
    fake = FakeListChatModel(responses=["The document shows a figure with some text."])

    reader = DocReader.from_llm(fake, files_dir=tmp_path, extract_images=True)
    result = await reader._arun("report.docx", "what is in it?")

    # Summary text is present AND the verbatim image path survived.
    assert "/doc_images/" in result
    assert "view_image" in result
