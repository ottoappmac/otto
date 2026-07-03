"""Shared loading, chunking, and ranking utilities for research tools.

Provides a single source of truth for:
- DuckDuckGo search via the ``ddgs`` package
- Fetching URLs as markdown (``aiohttp`` + ``markdownify``)
- Loading local files (per-format lightweight parsers)
- Recursive directory scanning
- Text chunking (``RecursiveCharacterTextSplitter``)
- BM25 keyword ranking (``rank_bm25``)

File parsing uses small, PyInstaller-friendly libraries instead of the
heavyweight ``unstructured`` stack:

- PDF: ``pypdf``
- DOCX: ``python-docx``
- PPTX: ``python-pptx``
- XLSX: ``openpyxl``
- HTML: ``markdownify``
- Plain text formats: stdlib
"""

from __future__ import annotations

import asyncio
import concurrent.futures
import functools
import hashlib
import io
import json
import logging
import os
import ssl
import threading
import time
from collections.abc import Callable
from dataclasses import dataclass
from pathlib import Path
from typing import TYPE_CHECKING, Optional, TypeAlias

if TYPE_CHECKING:
    from ddgs import DDGS

import aiohttp
import certifi
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

logger = logging.getLogger(__name__)

_MAX_CONTENT_CHARS = 80_000

# ---------------------------------------------------------------------------
# Image extraction (VLM-gated)
# ---------------------------------------------------------------------------

# Decorative artwork (bullets, logos, rules) tends to be tiny.  Skip anything
# below these thresholds so the agent isn't flooded with noise placeholders.
_MIN_IMAGE_BYTES = 1024
_MIN_IMAGE_DIM = 64

# OOXML namespaces used when walking DOCX bodies in document order.
_NS_W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


@dataclass
class _ImageSink:
    """Collects images extracted from a document into a session directory.

    ``images_dir`` is the absolute directory images are written under; each
    document gets its own ``{images_dir}/{doc_key}/`` subdir.  ``ref_prefix``
    is the *session-relative* path prefix the agent uses with ``view_image``
    (e.g. ``"/doc_images"`` when ``images_dir`` is ``{files_dir}/doc_images``).
    """

    images_dir: Path
    ref_prefix: str

    def add(
        self,
        doc_key: str,
        raw: bytes,
        ext: str,
        seq: int,
        *,
        page: int | None = None,
        slide: int | None = None,
    ) -> str | None:
        """Persist *raw* image bytes and return its session-relative ref.

        Returns ``None`` when the image is filtered out (too small) so the
        caller skips emitting a placeholder for it.
        """
        if not _keep_image(raw):
            return None
        doc_dir = self.images_dir / doc_key
        doc_dir.mkdir(parents=True, exist_ok=True)
        loc = page if page is not None else slide
        stem = f"p{loc}_img{seq}" if loc is not None else f"img{seq}"
        if not ext.startswith("."):
            ext = "." + ext if ext else ".png"
        name = f"{stem}{ext}"
        dest = doc_dir / name
        # mtime is folded into doc_key, so an identical-size file already on
        # disk is the same extraction — skip the rewrite.
        if not (dest.exists() and dest.stat().st_size == len(raw)):
            dest.write_bytes(raw)
        return f"{self.ref_prefix}/{doc_key}/{name}"

    def write_manifest(self, doc_key: str, source: str, entries: list[dict]) -> None:
        """Write a per-document ``manifest.json`` recording image order."""
        if not entries:
            return
        doc_dir = self.images_dir / doc_key
        doc_dir.mkdir(parents=True, exist_ok=True)
        (doc_dir / "manifest.json").write_text(
            json.dumps(
                {"source": source, "doc_key": doc_key, "images": entries},
                indent=2,
            ),
            encoding="utf-8",
        )


_FileParser: TypeAlias = Callable[[Path, Optional[_ImageSink]], list[Document]]


def _doc_key(path: Path) -> str:
    """Stable per-document key combining the file stem and a path+mtime hash.

    Folding mtime into the hash means editing a file produces a new key (and
    thus a fresh extraction), while re-loading an unchanged file reuses the
    already-written images.
    """
    try:
        mtime = path.stat().st_mtime_ns
    except OSError:
        mtime = 0
    digest = hashlib.sha1(f"{path.resolve()}:{mtime}".encode()).hexdigest()[:8]
    stem = "".join(c if c.isalnum() or c in "-_" else "_" for c in path.stem)[:48]
    return f"{stem}_{digest}"


def _keep_image(raw: bytes) -> bool:
    """Return True if *raw* looks like a meaningful (non-decorative) image."""
    if len(raw) < _MIN_IMAGE_BYTES:
        return False
    try:
        from PIL import Image

        with Image.open(io.BytesIO(raw)) as im:
            w, h = im.size
        return w >= _MIN_IMAGE_DIM and h >= _MIN_IMAGE_DIM
    except Exception:
        # Can't inspect dimensions — keep it rather than silently dropping.
        return True


def _image_placeholder(
    ref: str,
    seq: int,
    *,
    page: int | None = None,
    slide: int | None = None,
) -> str:
    """Build the inline placeholder line emitted in document text.

    Emitted on its own paragraph (callers surround it with blank lines) so the
    text splitter prefers to keep it attached to neighbouring text, preserving
    the "what comes before/after this image" adjacency.
    """
    if page is not None:
        loc = f"page {page}, "
    elif slide is not None:
        loc = f"slide {slide}, "
    else:
        loc = ""
    return (
        f"[IMAGE: {ref} — {loc}image {seq}. "
        f"Use view_image with this path to inspect it or ask a question about it.]"
    )

# Dedicated thread pool for heavy research I/O (DDG search, markdownify,
# PDF parsing) so these don't starve the default asyncio executor used by
# the backend's polling endpoints and file writes.
_RESEARCH_POOL = concurrent.futures.ThreadPoolExecutor(
    max_workers=8, thread_name_prefix="research",
)


async def _run_in_research_pool(fn, *args):
    """Run *fn* in the dedicated research thread pool."""
    try:
        pending = _RESEARCH_POOL._work_queue.qsize()
        active = len([t for t in _RESEARCH_POOL._threads if t.is_alive()])
        if pending > 0:
            logger.debug("[pool] research pool: %d queued, %d/%d threads active", pending, active, _RESEARCH_POOL._max_workers)
    except AttributeError:
        pass
    loop = asyncio.get_running_loop()
    return await loop.run_in_executor(_RESEARCH_POOL, fn, *args)


_SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx", ".pptx", ".xlsx",
    ".html", ".htm",
    ".txt", ".md", ".rst", ".json", ".xml", ".csv",
}


# ---------------------------------------------------------------------------
# DuckDuckGo search via ddgs package
# ---------------------------------------------------------------------------

_ddgs_local = threading.local()


def _get_ddgs() -> "DDGS":
    """Return a per-thread DDGS instance via ``threading.local()``.

    Each thread in the research pool gets its own ``DDGS`` (and therefore
    its own ``primp.Client``).  The Rust-backed ``primp.Client`` deadlocks
    when the same instance is called concurrently from multiple threads,
    so a shared singleton is unsafe.  Per-thread instances avoid this
    while still reusing the expensive TLS setup across searches within
    the same thread.

    Note: ``_get_ddg_semaphore()`` already serialises DDG searches at the
    async layer, so concurrency is the primary defence.  Per-thread
    instances are defence-in-depth in case the semaphore is ever relaxed.
    """
    inst = getattr(_ddgs_local, "instance", None)
    if inst is None:
        t0 = time.monotonic()
        from ddgs import DDGS
        inst = DDGS(timeout=20)
        _ddgs_local.instance = inst
        logger.debug("[ddgs] DDGS instance created for thread %s in %.1fs",
                     threading.current_thread().name, time.monotonic() - t0)
    return inst


def _warm_ddgs() -> None:
    """Pre-initialise a DDGS instance on this thread so the import and
    Rust TLS setup happen before real traffic arrives.
    """
    try:
        ddgs = _get_ddgs()
        ddgs.text("ping", max_results=1)
        logger.info("DDGS pre-warmed (thread=%s)", threading.current_thread().name)
    except Exception:
        logger.debug("DDGS warm-up failed (will retry on first search)", exc_info=True)


def warm_up() -> None:
    """Kick off background warm-up of heavy singletons (DDGS, SSL ctx).

    Call once from the application lifespan/startup rather than at import
    time so that tests and scripts that import this module don't pay for
    unwanted side-effects.
    """
    threading.Thread(target=_warm_ddgs, daemon=True, name="ddgs-warmup").start()
    _get_ssl_ctx()


_ddg_semaphore: asyncio.Semaphore | None = None


def _get_ddg_semaphore() -> asyncio.Semaphore:
    """Lazily create a per-loop semaphore that serialises DDG searches.

    ``primp`` (the Rust HTTP client inside ``ddgs``) deadlocks when
    multiple calls run concurrently in the same process — even on
    separate ``primp.Client`` instances.  Serialising at the async level
    keeps the event loop responsive while avoiding the Rust-level
    deadlock.
    """
    global _ddg_semaphore
    if _ddg_semaphore is None:
        _ddg_semaphore = asyncio.Semaphore(1)
    return _ddg_semaphore


async def ddg_search_urls(
    query: str,
    max_results: int = 10,
) -> list[dict]:
    """Search DuckDuckGo and return result dicts.

    Returns a list of ``{"url": ..., "title": ...}`` dicts.
    Uses the ``ddgs`` package with worldwide region to avoid
    locale-biased results.

    Searches are serialised via semaphore because the Rust-backed
    ``primp.Client`` inside ``ddgs`` deadlocks under concurrent use
    in PyInstaller frozen binaries.
    """

    def _search() -> list[dict]:
        tid = threading.current_thread().name
        logger.debug("[ddg] thread=%s query=%r — calling ddgs.text()", tid, query)
        t0 = time.monotonic()
        ddgs = _get_ddgs()
        results = [
            {"url": r["href"], "title": r["title"], "snippet": r["body"]}
            for r in ddgs.text(query, region="wt-wt", max_results=max_results)
        ]
        logger.debug("[ddg] thread=%s query=%r — done in %.1fs (%d results)", tid, query, time.monotonic() - t0, len(results))
        return results

    t0 = time.monotonic()
    async with _get_ddg_semaphore():
        results = await _run_in_research_pool(_search)
    logger.info("[ddg] %d results for %r (%.1fs)", len(results), query, time.monotonic() - t0)
    return results


# ---------------------------------------------------------------------------
# URL fetching
# ---------------------------------------------------------------------------

_ssl_ctx: ssl.SSLContext | None = None


def _get_ssl_ctx() -> ssl.SSLContext:
    """Return a cached SSL context (avoids re-parsing the CA bundle on each call).

    Not protected by a lock — the worst case under a race is two threads
    both creating a context and one write wins; ``ssl.SSLContext`` is
    thread-safe once constructed, so this is harmless.
    """
    global _ssl_ctx
    if _ssl_ctx is None:
        _ssl_ctx = ssl.create_default_context(cafile=certifi.where())
    return _ssl_ctx


async def fetch_url_as_markdown(url: str, timeout: int = 15) -> str:
    """Fetch *url* and return the page content as markdown text.

    Uses ``markdownify`` to convert HTML.  PDF responses are parsed via
    ``pypdf``.  Heavy parsing runs in a thread to avoid blocking the
    async event loop.
    """
    from markdownify import markdownify

    t0 = time.monotonic()
    logger.debug("[fetch] connecting to %s", url)
    connector = aiohttp.TCPConnector(ssl=_get_ssl_ctx())
    ua = "Mozilla/5.0 (compatible; ResearchTool/1.0)"

    async with aiohttp.ClientSession(
        connector=connector,
        max_field_size=65_536,
    ) as session:
        async with session.get(
            url,
            headers={"User-Agent": ua},
            allow_redirects=True,
            timeout=aiohttp.ClientTimeout(total=timeout),
        ) as resp:
            resp.raise_for_status()
            content_type = resp.headers.get("Content-Type", "").lower()
            raw = await resp.read()
    logger.debug("[fetch] downloaded %s (%d bytes, %.1fs)", url, len(raw), time.monotonic() - t0)

    url_path = url.lower().split("?")[0]
    if "application/pdf" in content_type or url_path.endswith(".pdf"):
        text = await _run_in_research_pool(_parse_pdf_bytes, raw)
        return _truncate(text)

    html = raw.decode("utf-8", errors="replace")
    md = await _run_in_research_pool(markdownify, html)
    md = md.strip()
    if len(md) > _MAX_CONTENT_CHARS:
        md = md[:_MAX_CONTENT_CHARS] + "\n\n...(truncated)"
    logger.debug("[fetch] converted %s to markdown (%d chars, %.1fs total)", url, len(md), time.monotonic() - t0)
    return md


# ---------------------------------------------------------------------------
# Local file / directory loading
# ---------------------------------------------------------------------------

_PARSERS: dict[str, _FileParser] = {}  # populated by _register_parsers()


def _register_parsers() -> None:
    """Build the extension → parser dispatch table.

    Called once at module level.  Each parser is a callable
    ``(Path) -> list[Document]``.
    """
    _PARSERS.update({
        ".pdf": _parse_pdf_file,
        ".docx": _parse_docx_file,
        ".pptx": _parse_pptx_file,
        ".xlsx": _parse_xlsx_file,
        ".html": _parse_html_file,
        ".htm": _parse_html_file,
    })
    for ext in (".txt", ".md", ".rst", ".json", ".xml", ".csv"):
        _PARSERS[ext] = _parse_text_file


def load_file(path: Path, sink: _ImageSink | None = None) -> list[Document]:
    """Load a single local file using a lightweight per-format parser.

    When *sink* is provided, image-aware parsers (PDF/DOCX/PPTX) extract
    embedded images into it and emit inline placeholders.
    """
    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    ext = path.suffix.lower()
    parser = _PARSERS.get(ext)
    if parser is None:
        raise ValueError(f"Unsupported file type: {ext}")
    return parser(path, sink)


def load_directory(
    directory: Path,
    recursive: bool = True,
    sink: _ImageSink | None = None,
) -> list[Document]:
    """Load all supported files in *directory*."""
    pattern = "**/*" if recursive else "*"
    docs: list[Document] = []
    for p in directory.glob(pattern):
        if p.is_file() and p.suffix.lower() in _SUPPORTED_EXTENSIONS:
            try:
                docs.extend(load_file(p, sink))
            except Exception as exc:
                logger.warning("Skipping %s: %s", p, exc)
    logger.info("Loaded %d documents from directory: %s", len(docs), directory)
    return docs


def _build_sink(
    files_dir: Path | None,
    extract_images: bool,
    images_dir: Path | None,
) -> _ImageSink | None:
    """Construct an :class:`_ImageSink` when image extraction is requested.

    ``images_dir`` defaults to ``{files_dir}/doc_images`` so extracted images
    live inside the session tree and are reachable by ``view_image`` (which
    resolves paths relative to ``files_dir``).
    """
    if not extract_images:
        return None
    if images_dir is None:
        if files_dir is None:
            return None
        images_dir = Path(files_dir) / "doc_images"
    images_dir = Path(images_dir)
    return _ImageSink(images_dir=images_dir, ref_prefix="/" + images_dir.name)


async def load_source(
    source: str,
    files_dir: Path | None = None,
    *,
    extract_images: bool = False,
    images_dir: Path | None = None,
) -> list[Document]:
    """Route *source* to the right loader (URL, file, or directory).

    Parameters
    ----------
    source:
        A URL (``http://`` / ``https://``), or a file/directory path.
    files_dir:
        If provided and *source* looks like a session-relative path
        (starts with ``/``), the path is resolved under *files_dir*.
        This lets the ``DocReader`` tool locate uploaded files that the
        agent references as ``/uploads/foo.pdf``.
    extract_images:
        When True (and a vision model is available upstream), embedded
        images in PDF/DOCX/PPTX files are extracted to ``images_dir`` and
        inline placeholders are emitted in the returned text so the agent
        can ``view_image`` them.
    images_dir:
        Absolute directory to write extracted images to.  Defaults to
        ``{files_dir}/doc_images`` when omitted.
    """
    if source.startswith(("http://", "https://")):
        try:
            md = await fetch_url_as_markdown(source)
        except Exception as exc:
            logger.warning("Failed to fetch %s: %s", source, exc)
            md = f"[Error fetching {source}: {type(exc).__name__}: {exc}]"
        return [Document(page_content=md, metadata={"source": source})]

    p = Path(source)

    if files_dir is not None:
        # All local sources are sandboxed to the session files directory,
        # matching the containment enforced by view_image/file tools.
        # Containment is checked on the lexically-normalised path (not the
        # symlink-resolved one) because uploads and dropped folders are
        # symlinked *into* the session tree by design.
        base = Path(os.path.normpath(files_dir.absolute()))
        if p.is_absolute() and p.exists():
            if not Path(os.path.normpath(str(p))).is_relative_to(base):
                raise ValueError(f"Path is outside the session directory: {source}")
        else:
            # Agent passes "/uploads/foo.pdf" or "uploads/foo.pdf" — resolve
            # under the session files directory.
            p = Path(os.path.normpath(str(base / source.lstrip("/"))))
            if not p.is_relative_to(base):
                raise ValueError(f"Path is outside the session directory: {source}")

    sink = _build_sink(files_dir, extract_images, images_dir)

    try:
        if p.is_dir():
            return await _run_in_research_pool(
                functools.partial(load_directory, p, sink=sink),
            )
        return await _run_in_research_pool(
            functools.partial(load_file, p, sink),
        )
    except FileNotFoundError:
        logger.warning("File not found: %s (resolved to %s)", source, p)
        raise
    except Exception as exc:
        logger.warning("Failed to load %s: %s", source, exc)
        raise


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------

def split_documents(
    docs: list[Document],
    chunk_size: int = 1500,
    chunk_overlap: int = 150,
) -> list[Document]:
    """Split *docs* into smaller chunks."""
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
    )
    return splitter.split_documents(docs)


# ---------------------------------------------------------------------------
# BM25 ranking
# ---------------------------------------------------------------------------

def rank_chunks(
    chunks: list[Document],
    query: str,
    k: int = 10,
) -> list[Document]:
    """Return the top-*k* chunks ranked by BM25 relevance to *query*.

    Falls back to simple keyword overlap scoring when ``rank_bm25`` is
    unavailable (e.g. corrupted PyInstaller archive).
    """
    if not chunks:
        return []
    if len(chunks) <= k:
        return chunks

    try:
        from rank_bm25 import BM25Okapi

        tokenized = [c.page_content.lower().split() for c in chunks]
        bm25 = BM25Okapi(tokenized)
        scores = bm25.get_scores(query.lower().split())
    except Exception:
        query_tokens = set(query.lower().split())
        scores = [
            len(query_tokens & set(c.page_content.lower().split()))
            for c in chunks
        ]

    top_indices = sorted(
        range(len(scores)), key=lambda i: scores[i], reverse=True,
    )[:k]
    return [chunks[i] for i in top_indices]


# ---------------------------------------------------------------------------
# Per-format parsers
# ---------------------------------------------------------------------------

def _parse_pdf_bytes(raw: bytes) -> str:
    """Extract text from raw PDF bytes via ``pypdf``."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(raw))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(p for p in pages if p.strip())


def _parse_pdf_file(path: Path, sink: _ImageSink | None = None) -> list[Document]:
    """Parse a local PDF file, optionally extracting embedded images.

    Text is extracted per page.  When *sink* is set, each page's embedded
    images are saved and a placeholder appended after that page's text
    (pypdf cannot position an image within the page's text run, so ordering
    is page-level rather than within-page).
    """
    if sink is None:
        text = _parse_pdf_bytes(path.read_bytes())
        return [Document(page_content=text, metadata={"source": str(path)})]

    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(path.read_bytes()))
    doc_key = _doc_key(path)
    manifest: list[dict] = []
    page_blocks: list[str] = []

    for pageno, page in enumerate(reader.pages, 1):
        segments: list[str] = []
        text = page.extract_text() or ""
        if text.strip():
            segments.append(text.strip())

        seq = 0
        for img in _safe_pdf_images(page):
            seq += 1
            raw = getattr(img, "data", None)
            if not raw:
                continue
            ext = Path(getattr(img, "name", "") or "").suffix or ".png"
            ref = sink.add(doc_key, raw, ext, seq, page=pageno)
            if ref is None:
                continue
            segments.append(_image_placeholder(ref, seq, page=pageno))
            manifest.append({
                "image_id": Path(ref).name,
                "path": ref,
                "page": pageno,
                "seq_index": seq,
            })

        if segments:
            page_blocks.append("\n\n".join(segments))

    sink.write_manifest(doc_key, str(path), manifest)
    text = "\n\n".join(page_blocks)
    return [Document(page_content=text, metadata={"source": str(path)})]


def _safe_pdf_images(page) -> list:
    """Return a page's embedded images, tolerating pypdf decode errors."""
    try:
        return list(page.images)
    except Exception as exc:
        logger.debug("pypdf image extraction failed on a page: %s", exc)
        return []


def _parse_docx_file(path: Path, sink: _ImageSink | None = None) -> list[Document]:
    """Parse a ``.docx`` file via ``python-docx``.

    When *sink* is set, the document body is walked in order so inline image
    placeholders are emitted at the exact paragraph position where the image
    appears, preserving before/after text adjacency.
    """
    from docx import Document as DocxDocument

    doc = DocxDocument(str(path))

    if sink is None:
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        text = "\n\n".join(paragraphs)
        return [Document(page_content=text, metadata={"source": str(path)})]

    doc_key = _doc_key(path)
    manifest: list[dict] = []
    parts: list[str] = []
    seq = 0
    rel_parts = doc.part.related_parts
    blip_tag = f"{{{_NS_A}}}blip"
    embed_attr = f"{{{_NS_R}}}embed"
    t_tag = f"{{{_NS_W}}}t"

    for child in doc.element.body.iterchildren():
        tag = child.tag
        if tag == f"{{{_NS_W}}}p":
            para_text = "".join(t.text or "" for t in child.iter(t_tag)).strip()
            if para_text:
                parts.append(para_text)
            for blip in child.iter(blip_tag):
                rid = blip.get(embed_attr)
                if not rid or rid not in rel_parts:
                    continue
                part = rel_parts[rid]
                raw = getattr(part, "blob", None)
                if not raw:
                    continue
                seq += 1
                ext = Path(str(getattr(part, "partname", ""))).suffix or ".png"
                ref = sink.add(doc_key, raw, ext, seq)
                if ref is None:
                    continue
                parts.append(_image_placeholder(ref, seq))
                manifest.append({
                    "image_id": Path(ref).name,
                    "path": ref,
                    "seq_index": seq,
                })
        elif tag == f"{{{_NS_W}}}tbl":
            for row in child.iter(f"{{{_NS_W}}}tr"):
                cells = [
                    "".join(t.text or "" for t in cell.iter(t_tag)).strip()
                    for cell in row.iter(f"{{{_NS_W}}}tc")
                ]
                cells = [c for c in cells if c]
                if cells:
                    parts.append(" | ".join(cells))

    sink.write_manifest(doc_key, str(path), manifest)
    text = "\n\n".join(parts)
    return [Document(page_content=text, metadata={"source": str(path)})]


def _parse_pptx_file(path: Path, sink: _ImageSink | None = None) -> list[Document]:
    """Parse a ``.pptx`` file via ``python-pptx``.

    When *sink* is set, picture shapes are extracted and placeholders emitted
    within each slide block in shape order.
    """
    from pptx import Presentation

    prs = Presentation(str(path))
    doc_key = _doc_key(path) if sink is not None else None
    manifest: list[dict] = []
    slides: list[str] = []

    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        _picture_type = MSO_SHAPE_TYPE.PICTURE
    except Exception:
        _picture_type = None

    for i, slide in enumerate(prs.slides, 1):
        parts = [f"--- Slide {i} ---"]
        seq = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
            if (
                sink is not None
                and _picture_type is not None
                and getattr(shape, "shape_type", None) == _picture_type
            ):
                try:
                    image = shape.image
                    raw = image.blob
                    ext = "." + (image.ext or "png")
                except Exception:
                    continue
                seq += 1
                ref = sink.add(doc_key, raw, ext, seq, slide=i)
                if ref is None:
                    continue
                parts.append(_image_placeholder(ref, seq, slide=i))
                manifest.append({
                    "image_id": Path(ref).name,
                    "path": ref,
                    "slide": i,
                    "seq_index": seq,
                })
        if len(parts) > 1:
            slides.append("\n".join(parts))

    if sink is not None:
        sink.write_manifest(doc_key, str(path), manifest)
    text = "\n\n".join(slides)
    return [Document(page_content=text, metadata={"source": str(path)})]


def _parse_xlsx_file(path: Path, sink: _ImageSink | None = None) -> list[Document]:
    """Parse a ``.xlsx`` file via ``openpyxl``."""
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    sheets: list[str] = []
    for ws in wb.worksheets:
        rows: list[str] = [f"--- Sheet: {ws.title} ---"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if len(rows) > 1:
            sheets.append("\n".join(rows))
    wb.close()

    text = "\n\n".join(sheets)
    return [Document(page_content=text, metadata={"source": str(path)})]


def _parse_html_file(path: Path, sink: _ImageSink | None = None) -> list[Document]:
    """Parse a local HTML file via ``markdownify``."""
    from markdownify import markdownify

    html = path.read_text(encoding="utf-8", errors="replace")
    text = markdownify(html).strip()
    return [Document(page_content=text, metadata={"source": str(path)})]


def _parse_text_file(path: Path, sink: _ImageSink | None = None) -> list[Document]:
    """Read a plain-text file (txt, md, csv, json, xml, rst)."""
    text = path.read_text(encoding="utf-8", errors="replace")
    return [Document(page_content=text, metadata={"source": str(path)})]


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _truncate(text: str) -> str:
    """Truncate *text* to ``_MAX_CONTENT_CHARS``."""
    text = text.strip()
    if len(text) > _MAX_CONTENT_CHARS:
        return text[:_MAX_CONTENT_CHARS] + "\n\n...(truncated)"
    return text


# ---------------------------------------------------------------------------
# Initialise parser dispatch table (must come after parser definitions)
# ---------------------------------------------------------------------------

_register_parsers()
